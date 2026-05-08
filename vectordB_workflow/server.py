import os, re, uuid, asyncio, logging, io, base64
from contextlib import asynccontextmanager
from typing import Optional, Literal, List
from datetime import datetime, timedelta
from pathlib import Path

import uvicorn, asyncpg
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from dotenv import load_dotenv

# LangChain & AI
from langchain_openai import AzureOpenAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import WebBaseLoader
from langchain_core.documents import Document
from langchain_core.embeddings import Embeddings

# Document Parsers (In-Memory)
from pypdf import PdfReader
import docx2txt
from pptx import Presentation

# Azure Storage
from azure.storage.blob import (
    generate_blob_sas,
    BlobSasPermissions,
    BlobServiceClient,
    CorsRule
)

load_dotenv()

# Logging Configuration
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s", datefmt="%Y-%m-%d %H:%M:%S")
logger = logging.getLogger("UniversalDocAPI")
os.environ["USER_AGENT"] = "Universal Document API 2.0"

# ==============================================================================
# CONFIGURATION
# ==============================================================================
class Config:
    """Application configuration from environment variables."""
    AZURE_STORAGE_ACCOUNT = os.getenv("AZURE_STORAGE_ACCOUNT")
    AZURE_STORAGE_KEY = os.getenv("AZURE_STORAGE_KEY")
    AZURE_CONTAINER = os.getenv("AZURE_CONTAINER")
    
    AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
    AZURE_OPENAI_API_INSTANCE_NAME = os.getenv("AZURE_OPENAI_API_INSTANCE_NAME")
    AZURE_OPENAI_EMBEDDING_DEPLOYMENT = os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT")
    AZURE_OPENAI_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION")
    AZURE_OPENAI_EMBEDDING_MODEL = os.getenv("AZURE_OPENAI_EMBEDDING_MODEL")
    
    POSTGRES_USER = os.getenv("POSTGRES_USER")
    POSTGRES_PASSWORD = os.getenv("POSTGRES_PASSWORD")
    POSTGRES_HOST = os.getenv("POSTGRES_HOST")
    POSTGRES_PORT = os.getenv("POSTGRES_PORT")
    POSTGRES_DATABASE = os.getenv("POSTGRES_DATABASE")
    
    CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", 500))
    CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", 50))
    MAX_FILE_SIZE_MB = int(os.getenv("MAX_FILE_SIZE_MB", 20))
    EMBEDDING_BATCH_SIZE = int(os.getenv("EMBEDDING_BATCH_SIZE", 100))
    ALLOWED_CORS_ORIGINS = os.getenv("ALLOWED_CORS_ORIGINS", "*").split(",")

config = Config()

# ==============================================================================
# PYDANTIC MODELS
# ==============================================================================
class UploadSASRequest(BaseModel):
    college_id: str
    document_id: str
    file_name: str

class ReadSASRequest(BaseModel):
    blob_path: str

class InsertDocumentRequest(BaseModel):
    document_name: Optional[str] = Field(None, min_length=1, max_length=255)
    document_url: Optional[str] = Field(None, max_length=2000)
    source_type: Literal["base64", "website"]
    source: Literal["Product-AI", "Zox-edu-ai", "zox-internal"] = Field("Product-AI")
    document_content: Optional[str] = None
    document_type: Optional[Literal["pdf", "docx", "pptx", "audio"]] = None
    chunk_size: Optional[int] = Field(None, ge=100, le=5000)
    chunk_overlap: Optional[int] = Field(None, ge=0, le=1000)
    trial_id: Optional[str] = Field(None, description="College Trial UUID")

class InsertFromBlobRequest(BaseModel):
    document_id: str = Field(..., description="Dataverse record GUID")
    blob_path: str = Field(..., description="Blob path from /blob/upload-sas response")
    document_name: str = Field(..., min_length=1, max_length=255)
    document_type: Literal["pdf", "docx", "pptx", "audio"] = Field(..., description="File type")
    source: Literal["Product-AI", "Zox-edu-ai", "zox-internal"] = Field("Product-AI")
    chunk_size: Optional[int] = Field(None, ge=100, le=5000)
    chunk_overlap: Optional[int] = Field(None, ge=0, le=1000)
    trial_id: Optional[str] = Field(None, description="College Trial UUID")

class RetrievalRequest(BaseModel):
    query: str = Field(..., min_length=1, max_length=500)
    limit: int = Field(5, ge=1, le=100)
    source: Literal["Product-AI", "Zox-edu-ai"] = Field("Product-AI")
    document_type: Optional[str] = None
    job_id: Optional[str] = None
    trial_id: Optional[str] = Field(None, description="College Trial UUID")

# ==============================================================================
# DATABASE CLIENT
# ==============================================================================
class PostgresClient:
    """PostgreSQL connection pool with lazy initialization."""
    def __init__(self):
        self._pool: Optional[asyncpg.Pool] = None
        self._lock = asyncio.Lock()

    async def _get_pool(self) -> asyncpg.Pool:
        if self._pool: return self._pool
        async with self._lock:
            if self._pool is None:
                logger.info("Initializing PostgreSQL connection pool...")
                self._pool = await asyncpg.create_pool(
                    user=config.POSTGRES_USER, password=config.POSTGRES_PASSWORD,
                    host=config.POSTGRES_HOST, port=config.POSTGRES_PORT,
                    database=config.POSTGRES_DATABASE, min_size=1, max_size=5,
                    command_timeout=60
                )
            return self._pool

    async def executemany(self, query: str, args):
        pool = await self._get_pool()
        async with pool.acquire() as conn:
            return await conn.executemany(query, args)

    async def execute(self, query: str, *args):
        pool = await self._get_pool()
        async with pool.acquire() as conn:
            return await conn.execute(query, *args)

    async def fetch(self, query: str, *args):
        pool = await self._get_pool()
        async with pool.acquire() as conn:
            return await conn.fetch(query, *args)

    async def fetchrow(self, query: str, *args):
        pool = await self._get_pool()
        async with pool.acquire() as conn:
            return await conn.fetchrow(query, *args)

    async def fetchval(self, query: str, *args):
        pool = await self._get_pool()
        async with pool.acquire() as conn:
            return await conn.fetchval(query, *args)

    async def close(self):
        if self._pool:
            await self._pool.close()
            self._pool = None

db_pool = PostgresClient()

# ==============================================================================
# AI & DOCUMENT TOOLS
# ==============================================================================
embeddings_model = AzureOpenAIEmbeddings(
    api_key=config.AZURE_OPENAI_API_KEY,
    azure_endpoint=config.AZURE_OPENAI_API_INSTANCE_NAME,
    azure_deployment=config.AZURE_OPENAI_EMBEDDING_DEPLOYMENT,
    api_version=config.AZURE_OPENAI_API_VERSION,
    model=config.AZURE_OPENAI_EMBEDDING_MODEL,
)

SUPPORTED_TYPES = {"pdf", "docx", "pptx"}

def download_blob_to_bytes(blob_path: str) -> bytes:
    """Download blob from Azure Storage using account key."""
    blob_name = blob_path.replace(f"{config.AZURE_CONTAINER}/", "", 1)
    connection_string = (
        f"DefaultEndpointsProtocol=https;"
        f"AccountName={config.AZURE_STORAGE_ACCOUNT};"
        f"AccountKey={config.AZURE_STORAGE_KEY};"
        f"EndpointSuffix=core.windows.net"
    )
    blob_service = BlobServiceClient.from_connection_string(connection_string)
    blob_client = blob_service.get_blob_client(container=config.AZURE_CONTAINER, blob=blob_name)
    return blob_client.download_blob().readall()

def extract_text_from_bytes(file_bytes: bytes, document_type: str) -> str:
    """Extract text from raw bytes (no temp files, 100% in-memory)."""
    doc_type = document_type.lower()
    stream = io.BytesIO(file_bytes)
    try:
        if doc_type == "pdf":
            reader = PdfReader(stream)
            text = "\n".join(page.extract_text() or "" for page in reader.pages).strip()
        elif doc_type == "docx":
            text = docx2txt.process(stream).strip()
        elif doc_type == "pptx":
            prs = Presentation(stream)
            parts = [shape.text_frame.text for slide in prs.slides for shape in slide.shapes if shape.has_text_frame]
            text = "\n".join(parts).strip()
        else:
            raise ValueError(f"Unsupported type: {document_type}")
        # Remove null bytes — PostgreSQL rejects \x00 in UTF-8 text columns
        # Common in scanned PDFs and PDFs with embedded fonts
        text = text.replace("\x00", "")
        return text
    finally:
        stream.close()

async def extract_website_data(url: str) -> str:
    """Extract website content using LangChain loader."""
    docs = await asyncio.to_thread(WebBaseLoader(url).load)
    return "\n".join(doc.page_content for doc in docs).strip()

async def embed_documents_batch(texts: List[str]) -> List[List[float]]:
    """Embed documents in batches."""
    batch_size = config.EMBEDDING_BATCH_SIZE
    result = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        result.extend(await asyncio.to_thread(embeddings_model.embed_documents, batch))
    return result

# ==============================================================================
# FASTAPI APPLICATION
# ==============================================================================
async def ensure_azure_cors():
    """Ensure Azure Blob Storage CORS rules are set for Dynamics 365."""
    try:
        if not config.AZURE_STORAGE_ACCOUNT or not config.AZURE_STORAGE_KEY:
            return
        connection_string = f"DefaultEndpointsProtocol=https;AccountName={config.AZURE_STORAGE_ACCOUNT};AccountKey={config.AZURE_STORAGE_KEY};EndpointSuffix=core.windows.net"
        blob_service_client = BlobServiceClient.from_connection_string(connection_string)
        cors_rule = CorsRule(
            allowed_origins=config.ALLOWED_CORS_ORIGINS,
            allowed_methods=["GET", "PUT", "POST", "OPTIONS", "HEAD", "MERGE", "DELETE"],
            allowed_headers=["*"],
            exposed_headers=["*"],
            max_age_in_seconds=3600
        )
        blob_service_client.set_service_properties(cors=[cors_rule])
        logger.info("✓ Azure Storage CORS rules validated.")
    except Exception as e:
        logger.error(f"Failed to configure Azure Storage CORS: {e}")

@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info("Universal Document API starting...")
    await ensure_azure_cors()
    yield
    await db_pool.close()

app = FastAPI(title="Universal Document API", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=config.ALLOWED_CORS_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ==============================================================================
# API ROUTES
# ==============================================================================
@app.get("/health")
async def health_check():
    result = await db_pool.fetchrow("SELECT NOW() as now")
    return {"status": "healthy", "database": "connected", "timestamp": str(result["now"])}

@app.post("/blob/upload-sas")
async def generate_upload_sas(request: UploadSASRequest):
    try:
        clean_filename = re.sub(r'[^a-zA-Z0-9.-]', '_', request.file_name)
        blob_name = f"{request.college_id}/{request.document_id}_{clean_filename}"
        
        now = datetime.utcnow()
        sas_token = generate_blob_sas(
            account_name=config.AZURE_STORAGE_ACCOUNT,
            container_name=config.AZURE_CONTAINER,
            blob_name=blob_name,
            account_key=config.AZURE_STORAGE_KEY,
            permission=BlobSasPermissions(write=True, create=True),
            start=now - timedelta(minutes=2),
            expiry=now + timedelta(minutes=15)
        )
        return {
            "upload_url": f"https://{config.AZURE_STORAGE_ACCOUNT}.blob.core.windows.net/{config.AZURE_CONTAINER}/{blob_name}?{sas_token}",
            "blob_path": f"{config.AZURE_CONTAINER}/{blob_name}"
        }
    except Exception as e:
        logger.error(f"SAS generation failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/blob/read-sas")
async def generate_read_sas(request: ReadSASRequest):
    """Generate a short-lived SAS token for reading a blob."""
    try:
        blob_name = request.blob_path.replace(f"{config.AZURE_CONTAINER}/", "")
        now = datetime.utcnow()
        sas_token = generate_blob_sas(
            account_name=config.AZURE_STORAGE_ACCOUNT,
            container_name=config.AZURE_CONTAINER,
            blob_name=blob_name,
            account_key=config.AZURE_STORAGE_KEY,
            permission=BlobSasPermissions(read=True),
            expiry=now + timedelta(minutes=15)
        )
        return {
            "read_url": f"https://{config.AZURE_STORAGE_ACCOUNT}.blob.core.windows.net/{config.AZURE_CONTAINER}/{blob_name}?{sas_token}"
        }
    except Exception as e:
        logger.error(f"Read SAS generation failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/insert")
async def insert_document(request: InsertDocumentRequest):
    """Legacy base64/website insertion with trial_id support."""
    job_id = uuid.uuid4()
    try:
        if request.source_type == "base64":
            content_bytes = base64.b64decode(request.document_content)
            text = extract_text_from_bytes(content_bytes, request.document_type)
            doc_type = request.document_type.lower()
        elif request.source_type == "website":
            text = await extract_website_data(request.document_url or request.document_content)
            doc_type = "website"
        else:
            raise HTTPException(status_code=400, detail="Invalid source type")

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=request.chunk_size or config.CHUNK_SIZE,
            chunk_overlap=request.chunk_overlap or config.CHUNK_OVERLAP
        )
        chunks = splitter.split_text(text)
        embeddings_list = await embed_documents_batch(chunks)

        await db_pool.executemany(
            """INSERT INTO education_vector_documents 
               (id, document_name, embedding, content, job_id, document_type, 
                chunk_index, source, trial_id)
               VALUES ($1, $2, $3::vector, $4, $5, $6::document_type_enum, $7, $8::source_enum, $9::uuid)""",
            [
                (uuid.uuid4(), request.document_name, str(e), chunk, str(job_id), 
                 doc_type, i, request.source, request.trial_id)
                for i, (chunk, e) in enumerate(zip(chunks, embeddings_list))
            ]
        )
        return {"success": True, "job_id": str(job_id), "total_chunks": len(chunks)}
    except Exception as e:
        logger.error(f"Insert failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/insert-from-blob")
async def insert_from_blob(request: InsertFromBlobRequest):
    job_id = uuid.uuid4()
    clean_filename = re.sub(r'[^a-zA-Z0-9.-]', '_', request.document_name)
    doc_identifier = f"{request.document_id}_{clean_filename}"

    try:
        logger.info(f"Processing blob job {job_id}: {request.document_name} (type: {request.document_type})")

        # Skip vector ingestion for audio files as requested
        if request.document_type == "audio":
            logger.info(f"⏭️ Skipping vector ingestion for audio file: {request.document_name}")
            return {
                "success": True,
                "job_id": str(job_id),
                "document_id": doc_identifier,
                "message": "Audio file recognized; skipped vector ingestion as per configuration."
            }

        # 1. Download
        blob_bytes = await asyncio.to_thread(download_blob_to_bytes, request.blob_path)
        
        # 2. Extract & Chunk
        text = extract_text_from_bytes(blob_bytes, request.document_type)
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=request.chunk_size or config.CHUNK_SIZE,
            chunk_overlap=request.chunk_overlap or config.CHUNK_OVERLAP
        )
        chunks = splitter.split_text(text)
        
        # 3. Embed
        embeddings_list = await embed_documents_batch(chunks)
        
        # 4. Transactional Upsert (Delete existing chunks first)
        pool = await db_pool._get_pool()
        async with pool.acquire() as conn:
            async with conn.transaction():
                await conn.execute("DELETE FROM education_vector_documents WHERE document_id = $1", doc_identifier)
                # Insert new chunks, storing the blob_path in document_url
                await conn.executemany(
                    """INSERT INTO education_vector_documents 
                       (id, document_name, document_url, embedding, content, job_id, 
                        document_type, chunk_index, source, document_id, trial_id)
                       VALUES ($1, $2, $3, $4::vector, $5, $6, $7::document_type_enum, 
                               $8, $9::source_enum, $10, $11::uuid)""",
                    [
                        (
                            uuid.uuid4(), request.document_name, request.blob_path, str(e), 
                            chunk, str(job_id), request.document_type.lower(), i, 
                            request.source, doc_identifier, request.trial_id
                        )
                        for i, (chunk, e) in enumerate(zip(chunks, embeddings_list))
                    ]
                )

        return {"success": True, "job_id": str(job_id), "document_id": doc_identifier, "total_chunks": len(chunks)}
    except Exception as e:
        logger.error(f"Insert failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/retrieve")
async def retrieve_documents(request: RetrievalRequest):
    try:
        query_emb = await asyncio.to_thread(embeddings_model.embed_query, request.query)
        query_str = "[" + ",".join(map(str, query_emb)) + "]"
        
        # Build query with optional filters
        filters = ["source = $3::source_enum"]
        params = [query_str, request.limit, request.source]
        idx = 4

        if request.trial_id:
            filters.append(f"trial_id = ${idx}::uuid")
            params.append(request.trial_id)
            idx += 1
        if request.document_type:
            filters.append(f"document_type = ${idx}::document_type_enum")
            params.append(request.document_type)
            idx += 1
        if request.job_id:
            filters.append(f"job_id = ${idx}::uuid")
            params.append(request.job_id)

        where_clause = f"WHERE {' AND '.join(filters)}"
        
        sql = f"""SELECT id, document_name, document_url, document_type::text, content, job_id, chunk_index, source,
                  1 - (embedding <=> $1::vector)::float as similarity_score
                  FROM education_vector_documents {where_clause} 
                  ORDER BY embedding <=> $1::vector LIMIT $2"""

        results = await db_pool.fetch(sql, *params)
        return {"success": True, "documents": [dict(r) for r in results]}
    except Exception as e:
        logger.error(f"Retrieval failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.delete("/document/{document_id}")
async def delete_document(document_id: str):
    """
    Delete a document from both the vector database and Azure Blob Storage.
    Applies sanitization to ensure names like '(1)' match the stored '__1_'.
    """
    try:
        # 1. Sanitize the incoming ID to match the stored version
        # This handles cases where '(1)' was stored as '__1_'
        sanitized_id = re.sub(r'[^a-zA-Z0-9.-]', '_', document_id)
        sql_pattern = f"{sanitized_id}%"
        
        # 2. Fetch blob paths from DB using sanitized ID
        find_sql = "SELECT DISTINCT document_url FROM education_vector_documents WHERE document_id LIKE $1"
        records = await db_pool.fetch(find_sql, sql_pattern)
        
        # 2. Setup Azure client
        connection_string = (
            f"DefaultEndpointsProtocol=https;"
            f"AccountName={config.AZURE_STORAGE_ACCOUNT};"
            f"AccountKey={config.AZURE_STORAGE_KEY};"
            f"EndpointSuffix=core.windows.net"
        )
        blob_service = BlobServiceClient.from_connection_string(connection_string)
        container_client = blob_service.get_container_client(config.AZURE_CONTAINER)
        
        deleted_blobs = []
        blob_paths_to_delete = set()

        # Add paths from DB
        for row in records:
            if row["document_url"]:
                blob_paths_to_delete.add(row["document_url"])

        # 3. DISCOVERY: If no paths in DB, or to be safe, search Azure directly for this GUID
        # This catches Audio files or orphaned files not in KB
        try:
            # We list blobs and look for the document_id in the name
            # Note: This is a best-effort scan.
            blobs_list = await asyncio.to_thread(container_client.list_blobs)
            for blob in blobs_list:
                if document_id in blob.name:
                    blob_paths_to_delete.add(blob.name)
        except Exception as scan_err:
            logger.warning(f"⚠️ Azure discovery scan failed: {scan_err}")

        # 4. Perform Deletion in Azure
        for blob_name in blob_paths_to_delete:
            try:
                # Normalize name (remove container prefix if stored in DB path)
                clean_name = blob_name.replace(f"{config.AZURE_CONTAINER}/", "", 1)
                await asyncio.to_thread(container_client.delete_blob, clean_name)
                deleted_blobs.append(clean_name)
                logger.info(f"🗑️ Deleted from Azure: {clean_name}")
            except Exception as b_err:
                logger.warning(f"⚠️ Could not delete blob {blob_name}: {b_err}")

        # 5. Delete chunks from PostgreSQL
        await db_pool.execute("DELETE FROM education_vector_documents WHERE document_id LIKE $1", sql_pattern)
        
        if not deleted_blobs and not records:
            raise HTTPException(status_code=404, detail=f"No document or storage file found for ID: {document_id}")

        logger.info(f"✅ Cleanup complete for {document_id}: {len(deleted_blobs)} files removed.")
        return {
            "success": True, 
            "document_id": document_id, 
            "blobs_deleted": deleted_blobs,
            "message": f"Purged {len(deleted_blobs)} file(s) from storage and cleared vector cache."
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Delete failed for '{document_id}': {e}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    uvicorn.run("server:app", host="0.0.0.0", port=8002, reload=True)
